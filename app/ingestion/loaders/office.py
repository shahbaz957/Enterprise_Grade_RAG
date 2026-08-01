"""Office loaders (.docx, .pptx)."""

from __future__ import annotations

from pathlib import Path

import logfire
from docx import Document
from pptx import Presentation

from app.ingestion.loaders.base import LoadedDocument, ensure_logfire, resolve_path

DOCX_EXTENSIONS = {".docx"}
PPTX_EXTENSIONS = {".pptx"}
OFFICE_EXTENSIONS = DOCX_EXTENSIONS | PPTX_EXTENSIONS


def load_docx(path: str | Path) -> LoadedDocument:
    ensure_logfire()
    file_path = resolve_path(path)

    with logfire.span("loader.docx", source=str(file_path)):
        document = Document(str(file_path))
        parts: list[str] = []

        for para in document.paragraphs:
            if para.text.strip():
                parts.append(para.text.strip())

        table_count = 0
        for table in document.tables:
            table_count += 1
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))

        doc = LoadedDocument(
            text="\n\n".join(parts),
            source=str(file_path),
            doc_type="docx",
            metadata={
                "filename": file_path.name,
                "extension": ".docx",
                "paragraph_count": len(document.paragraphs),
                "table_count": table_count,
                "size_bytes": file_path.stat().st_size,
            },
        )
        logfire.info(
            "Loaded DOCX document",
            source=str(file_path),
            chars=doc.char_count,
            tables=table_count,
        )
        return doc


def load_pptx(path: str | Path) -> LoadedDocument:
    ensure_logfire()
    file_path = resolve_path(path)

    with logfire.span("loader.pptx", source=str(file_path)):
        presentation = Presentation(str(file_path))
        slides_text: list[str] = []

        for idx, slide in enumerate(presentation.slides, start=1):
            chunks: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs).strip()
                    if line:
                        chunks.append(line)
            if chunks:
                slides_text.append(f"[Slide {idx}]\n" + "\n".join(chunks))

        doc = LoadedDocument(
            text="\n\n".join(slides_text),
            source=str(file_path),
            doc_type="pptx",
            metadata={
                "filename": file_path.name,
                "extension": ".pptx",
                "slide_count": len(presentation.slides),
                "size_bytes": file_path.stat().st_size,
            },
        )
        logfire.info(
            "Loaded PPTX document",
            source=str(file_path),
            slides=doc.metadata["slide_count"],
            chars=doc.char_count,
        )
        return doc


def load_office(path: str | Path) -> LoadedDocument:
    """Dispatch DOCX / PPTX by extension."""
    file_path = resolve_path(path)
    ext = file_path.suffix.lower()
    if ext in DOCX_EXTENSIONS:
        return load_docx(file_path)
    if ext in PPTX_EXTENSIONS:
        return load_pptx(file_path)
    raise ValueError(f"Unsupported office format: {ext} ({file_path})")
